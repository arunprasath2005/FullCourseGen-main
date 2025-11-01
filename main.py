import os
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, Field
from enum import Enum
import google.generativeai as genai
from dotenv import load_dotenv
import json
import uvicorn
import re
from fastapi.middleware.cors import CORSMiddleware
import requests
from docx import Document
from pdfplumber import open as pdf_open
from pptx import Presentation
from io import BytesIO
import asyncio

load_dotenv()

API_KEY = os.getenv("GOOGLE_GEMINI_KEY") or os.getenv("GOOGLE_API_KEY")
if not API_KEY:
    print("Warning: GOOGLE_GEMINI_KEY or GOOGLE_API_KEY not set in environment. Gemini requests may fail.")
genai.configure(api_key=API_KEY)

generation_config = {
    "temperature": 0.3,
    "top_p": 0.95,
    "top_k": 64,
    "max_output_tokens": 8192,
    "response_mime_type": "text/plain",
}

model = genai.GenerativeModel(
    model_name="gemini-2.0-flash-exp",
    generation_config=generation_config,
)

app = FastAPI(title="Full Course Generator API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

YOUTUBE_API_KEY = "AIzaSyAoV0_ILuFFE8WyfxbifJtk2asH2HFj9Do"

class DifficultyLevel(str, Enum):
    EASY = "easy"
    MEDIUM = "medium"
    HARD = "hard"

class UserLevel(str, Enum):
    BEGINNER = "Beginner"
    INTERMEDIATE = "Intermediate"
    ADVANCED = "Advanced"

class SubjectRequest(BaseModel):
    ques: str

class CourseRequest(BaseModel):
    subject: str = Field(..., description="The subject of the course")
    difficulty: DifficultyLevel = Field(..., description="Difficulty level of the course")
    focus_area: str = Field(..., description="Specific area to focus on within the subject")
    units: int = Field(..., ge=1, le=10, description="Number of units desired")

class FileRequest(BaseModel):
    file_url: str

class CourseRecommendationRequest(BaseModel):
    student_level: str
    course: str

class QuizResult(BaseModel):
    score: float = Field(..., ge=0, le=9, description="Score achieved in the quiz (0-9)")
    time_taken: float = Field(..., gt=0, description="Time taken to complete the quiz in seconds")

async def fetch_youtube_video(query: str) -> str:
    youtube_api_url = (
        f"https://www.googleapis.com/youtube/v3/search?part=snippet&maxResults=1&q={query}&type=video&key={YOUTUBE_API_KEY}"
    )
    try:
        response = requests.get(youtube_api_url, timeout=10)
        response.raise_for_status()
        data = response.json()
        if "items" in data and len(data["items"]) > 0:
            video_id = data["items"][0]["id"]["videoId"]
            return f"https://www.youtube.com/watch?v={video_id}"
        return "No relevant video found."
    except requests.exceptions.Timeout:
        return "YouTube fetch timeout."
    except Exception as e:
        print(f"Error fetching YouTube video: {e}")
        return "Error fetching video."

async def generate_unit_content(unit_data: dict, subject: str, difficulty: str, focus_area: str):
    content_prompt = f"""
    Generate detailed educational content for the unit "{unit_data['unitTitle']}" in {subject}.
    Topics to cover: {', '.join(unit_data['topicsCovered'])}
    Learning objectives: {', '.join(unit_data['learningObjectives'])}
    Difficulty level: {difficulty}
    Focus area: {focus_area}

    Return the response in this JSON format:
    {{
        "topicContents": [
            {{
                "topic": "Topic Name",
                "content": "Detailed explanation and educational content",
                "examples": ["example 1", "example 2"],
                "exercises": ["exercise 1", "exercise 2"]
            }}
        ]
    }}

    Ensure content is practical and matches the specified difficulty level.
    Give the content in about minimum 6000 words.
    """

    try:
        response = model.generate_content(content_prompt)
        cleaned_json = re.sub(r"^```json|```$", "", response.text, flags=re.MULTILINE).strip()
        content_data = json.loads(cleaned_json)
        return content_data
    except Exception as e:
        print(f"Error generating content for unit {unit_data['unitTitle']}: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate content for unit {unit_data['unitTitle']}: {str(e)}"
        )

async def get_unit_details(unit_title: str, subject: str, difficulty: str, focus_area: str):
    unit_prompt = f"""
    Generate a detailed unit structure for "{unit_title}" in {subject} course.
    Difficulty level: {difficulty}
    Focus area: {focus_area}

    Return the response in this JSON format:
    {{
        "unitTitle": "{unit_title}",
        "learningObjectives": ["detailed objective 1", "detailed objective 2"],
        "topicsCovered": ["detailed topic 1", "detailed topic 2"],
        "resources": ["resource 1", "resource 2"],
        "estimatedDuration": "X weeks"
    }}

    Ensure content matches the difficulty level and focuses on practical applications.
    """
    try:
        response = model.generate_content(unit_prompt)
        raw_response = response.text
        cleaned_json = re.sub(r"^```json|```$", "", raw_response, flags=re.MULTILINE).strip()

        try:
            unit_data = json.loads(cleaned_json)
        except json.JSONDecodeError as e:
            raise ValueError(f"JSON parsing error for {unit_title}: {str(e)}")

        detailed_content = await generate_unit_content(unit_data, subject, difficulty, focus_area)
        unit_data["detailedContent"] = detailed_content

        youtube_query = f"{unit_title} {subject} {focus_area}"
        youtube_video_url = await fetch_youtube_video(youtube_query)
        unit_data["youtube_video_url"] = youtube_video_url

        return unit_data
    except Exception as e:
        print(f"Error in get_unit_details for {unit_title}: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate unit details for {unit_title}: {str(e)}"
        )

async def generate_mcqs(unit_data: dict, subject: str, difficulty: str, focus_area: str):
    mcq_prompt = f"""
    Generate Multiple Choice Questions (MCQs) for the unit "{unit_data['unitTitle']}" in {subject}.
    Difficulty level: {difficulty}
    Focus area: {focus_area}

    Return the response in this JSON format:
    {{
        "unitAssessment": [
            {{
                "topic": "Topic Name",
                "questions": [
                    {{
                        "questionId": "unique_id",
                        "question": "Question text",
                        "options": [
                            "Option A",
                            "Option B",
                            "Option C",
                            "Option D"
                        ],
                        "correctAnswer": "Correct option",
                        "explanation": "Explanation of the correct answer"
                    }}
                ]
            }}
        ]
    }}

    Generate at least 3 MCQs per topic, and only 3 topics, ensuring they match the difficulty level.
    """

    try:
        response = model.generate_content(mcq_prompt)
        cleaned_json = re.sub(r"^```json|```$", "", response.text, flags=re.MULTILINE).strip()
        mcq_data = json.loads(cleaned_json)
        return mcq_data
    except Exception as e:
        print(f"Error generating MCQs for unit {unit_data['unitTitle']}: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate MCQs for unit {unit_data['unitTitle']}: {str(e)}"
        )

async def get_unit_details_with_mcqs(unit_title: str, subject: str, difficulty: str, focus_area: str):
    unit_prompt = f"""
    Generate a detailed unit structure for "{unit_title}" in {subject} course.
    Difficulty level: {difficulty}
    Focus area: {focus_area}

    Return the response in this JSON format:
    {{
        "unitTitle": "{unit_title}"
    }}

    Ensure content matches the difficulty level and focuses on practical applications.
    """

    try:
        response = model.generate_content(unit_prompt)
        cleaned_json = re.sub(r"^```json|```$", "", response.text, flags=re.MULTILINE).strip()
        unit_data = json.loads(cleaned_json)

        unit_mcqs = await generate_mcqs(unit_data, subject, difficulty, focus_area)
        unit_data["assessment"] = unit_mcqs

        return unit_data

    except Exception as e:
        print(f"Error in get_unit_details_with_mcqs for {unit_title}: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate unit details for {unit_title}: {str(e)}"
        )

def predict_user_level(score: float, time_taken: float) -> UserLevel:
    if score >= 7 and time_taken <= 80:
        return UserLevel.ADVANCED
    elif 4 <= score < 7:
        return UserLevel.INTERMEDIATE
    else:
        return UserLevel.BEGINNER

@app.get("/")
async def root():
    return {
        "message": "Full Course Generator API",
        "endpoints": [
            "/generate-course",
            "/doubt-chatbot",
            "/generate-question",
            "/predict-level",
            "/course-recommendation",
            "/detect-domain-from-file"
        ]
    }

@app.post("/doubt-chatbot")
async def doubt_chatbot(request: SubjectRequest):
    subject = request.ques

    prompt = f"""
    You are a doubt chatbot for students and you have to resolve students doubts. The question is: {subject}
    """

    response = model.generate_content(prompt)
    return {"answer": response.text}

@app.post("/generate-course")
async def generate_course(request: CourseRequest):
    try:
        structure_prompt = f"""
        Generate a comprehensive course structure for {request.subject} with exactly {request.units} units.
        Focus area: {request.focus_area}
        Difficulty: {request.difficulty}

        Return ONLY unit titles in this JSON format:
        {{
            "courseTitle": "",
            "difficultyLevel": "",
            "description": "",
            "prerequisites": ["prerequisite 1", "prerequisite 2"],
            "learningOutcomes": ["outcome 1", "outcome 2"],
            "units": [
                {{
                    "unitTitle": "",
                    "unitDescription": ""
                }}
            ],
            "overview": "",
            "assessmentMethods": ["method 1", "method 2"]
        }}
        """

        structure_response = model.generate_content(structure_prompt)
        cleaned_json = re.sub(r"^```json|```$", "", structure_response.text, flags=re.MULTILINE).strip()
        course_structure = json.loads(cleaned_json)

        detailed_units = []
        for unit in course_structure["units"]:
            try:
                unit_details = await get_unit_details(
                    unit["unitTitle"],
                    request.subject,
                    request.difficulty,
                    request.focus_area
                )
                detailed_units.append(unit_details)
            except Exception as e:
                print(f"Error processing unit {unit['unitTitle']}: {str(e)}")
                continue

        if not detailed_units:
            raise HTTPException(
                status_code=500,
                detail="Failed to generate any unit details"
            )

        course_structure["units"] = detailed_units
        return course_structure

    except Exception as e:
        print(f"Error in generate_course: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/generate-question")
async def generate_question(request: CourseRequest):
    try:
        structure_prompt = f"""
        Generate a comprehensive course structure for {request.subject} with exactly {request.units} units.
        Focus area: {request.focus_area}
        Difficulty: {request.difficulty}

        Return ONLY unit titles in this JSON format:
        {{
            "courseTitle": "",
            "difficultyLevel": "",
            "description": "",
            "prerequisites": ["prerequisite 1", "prerequisite 2"],
            "learningOutcomes": ["outcome 1", "outcome 2"],
            "units": [
                {{
                    "unitTitle": "",
                    "unitDescription": ""
                }}
            ],
            "overview": "",
            "assessmentMethods": ["method 1", "method 2"]
        }}
        """

        structure_response = model.generate_content(structure_prompt)
        cleaned_json = re.sub(r"^```json|```$", "", structure_response.text, flags=re.MULTILINE).strip()
        course_structure = json.loads(cleaned_json)

        unit_tasks = [
            get_unit_details_with_mcqs(
                unit["unitTitle"],
                request.subject,
                request.difficulty,
                request.focus_area
            ) for unit in course_structure["units"]
        ]

        detailed_units = await asyncio.gather(*unit_tasks, return_exceptions=True)

        course_structure["units"] = [
            unit for unit in detailed_units if not isinstance(unit, Exception)
        ]

        if not course_structure["units"]:
            raise HTTPException(
                status_code=500,
                detail="Failed to generate any unit details"
            )

        return course_structure

    except Exception as e:
        print(f"Error in generate_question: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/predict-level", response_model=UserLevel)
async def predict_level(quiz_result: QuizResult):
    try:
        level = predict_user_level(quiz_result.score, quiz_result.time_taken)
        return level
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Prediction error: {str(e)}")

@app.post("/course-recommendation")
async def recommend_course(request: CourseRecommendationRequest):
    student_level = request.student_level
    course = request.course

    prompt = f"""
    You are an intelligent assistant specializing in educational course recommendations.
    Based on the student's level and the specified course, recommend 4 appropriate courses with the following details:
    Subject, Number of Units, Focus Area, and Difficulty Level. Respond in JSON format.

    Input:
    1. Student Level: {student_level}
    2. Course: {course}

    Output:
        {{
            "subject": "Python",
            "units": 3,
            "focus_area": "Python Basics",
            "difficulty": "Beginner"
        }},
        {{
            "subject": "Data Structures",
            "units": 3,
            "focus_area": "Arrays and Linked Lists",
            "difficulty": "Intermediate"
        }},
        {{
            "subject": "Algorithms",
            "units": 3,
            "focus_area": "Sorting and Searching",
            "difficulty": "Intermediate"
        }},
        {{
            "subject": "Advanced Python",
            "units": 3,
            "focus_area": "Python for Data Science",
            "difficulty": "Advanced"
        }}
    """

    try:
        gemini_response = model.generate_content(prompt)
        response_text = gemini_response.text.strip()

        if response_text.startswith("```json"):
            response_text = response_text[7:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]
        if response_text.startswith("'''"):
            response_text = response_text[3:]
        if response_text.endswith("'''"):
            response_text = response_text[:-3]

        response_text = f"[{response_text.strip().strip('[]')}]"
        recommendations = json.loads(response_text)
        return {"recommendations": recommendations}

    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse recommendations: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred: {str(e)}")

@app.post("/detect-domain-from-file")
async def detect_domain_from_file(request: FileRequest):
    try:
        response = requests.get(request.file_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download file.")

        file_content = BytesIO(response.content)
        filename = os.path.basename(request.file_url)
        file_extension = filename.split(".")[-1].lower()

        content = ""
        if file_extension == "docx":
            doc = Document(file_content)
            content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        elif file_extension == "pdf":
            with pdf_open(file_content) as pdf:
                content = "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())
        elif file_extension == "pptx":
            ppt = Presentation(file_content)
            content = "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type.")

        if not content.strip():
            raise HTTPException(status_code=400, detail="Extracted content is empty.")

        prompt = f"""
        Analyze the following educational content and determine its subject domain (e.g., Mathematics, Physics, Biology, History, etc.)
        and subdomain (if applicable). Provide a brief explanation for why you classified it as that domain and subdomain.
        Format your response as JSON with three fields: 'domain', 'subdomain', and 'explanation'.

        Content: {content}
        """

        gemini_response = model.generate_content(prompt)
        response_text = gemini_response.text.strip()

        if response_text.startswith("```json"):
            response_text = response_text[7:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]
        if response_text.startswith("'''"):
            response_text = response_text[3:]
        if response_text.endswith("'''"):
            response_text = response_text[:-3]

        try:
            result = json.loads(response_text)
        except json.JSONDecodeError:
            raise HTTPException(status_code=500, detail="Gemini response is not in valid JSON format.")

        return {
            "filename": filename,
            "domain": result.get("domain", "Unknown"),
            "subdomain": result.get("subdomain", "Unknown"),
            "explanation": result.get("explanation", "No explanation provided."),
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
